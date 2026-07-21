"""SENTINEL Presentation Deck Generator
Creates a 12-slide PowerPoint for ET AI Hackathon 2.0
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# SENTINEL brand colors
DARK_BG = RGBColor(0x0E, 0x11, 0x17)
ACCENT_BLUE = RGBColor(0x00, 0xAA, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x9C, 0xA3, 0xAF)
DARK_GRAY = RGBColor(0x37, 0x41, 0x51)
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
YELLOW = RGBColor(0xFA, 0xCC, 0x15)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── SLIDE 1: TITLE ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, 1, 1.5, 11, 1.5, "SENTINEL", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.0, 11, 0.8, "Digital Public Safety Intelligence Platform", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.0, 11, 0.5, "ET AI Hackathon 2.0  |  AI for Digital Public Safety", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 5.5, 11, 0.5, "SCAMWatch  •  CURRENCYGuard  •  FRAUDGraph  •  GeoIntel", font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 2: THE PROBLEM ────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.8, "The Problem", font_size=36, color=WHITE, bold=True)

    problems = [
        "1.14 million cybercrime complaints in 2023 — up 60% from 2022",
        "Rs 1,776 crore lost to digital arrest scams in 9 months (2024)",
        "Counterfeit Rs 500 notes defeating manual detection at banks",
        "Reactive investigation — need predictive threat neutralisation",
        "Citizens lack real-time tools to verify suspicious messages",
    ]
    add_bullet_slide(slide, 0.8, 1.5, 11, 5, problems, font_size=20, color=WHITE)

    # ── SLIDE 3: OUR SOLUTION ───────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.8, "Our Solution: SENTINEL", font_size=36, color=WHITE, bold=True)

    solutions = [
        "4-module AI platform for proactive fraud detection",
        "Cross-module intelligence sharing via shared ChromaDB",
        "Real-time scam detection + citizen protection alerts",
        "Court-admissible evidence packaging for law enforcement",
        "Multi-language support (Hindi, Tamil, Bengali, Telugu)",
        "Geospatial threat mapping with 15 NCRB/RBI hotspots",
    ]
    add_bullet_slide(slide, 0.8, 1.5, 11, 5, solutions, font_size=20, color=WHITE)

    # ── SLIDE 4: ARCHITECTURE ───────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.8, "Architecture", font_size=36, color=WHITE, bold=True)

    arch_items = [
        "Frontend: Streamlit (multi-page, dark theme)",
        "Backend: FastAPI (11+ REST endpoints)",
        "LLM: Groq Llama 3.3 70B via langchain-groq",
        "Agent Framework: LangGraph (5/4/6-node pipelines)",
        "Vector Store: ChromaDB (shared intelligence layer)",
        "Computer Vision: OpenCV (7 security feature checks)",
        "Graph Analysis: NetworkX + pyvis visualization",
        "Graph Database: Neo4j Aura (with in-memory fallback)",
    ]
    add_bullet_slide(slide, 0.8, 1.5, 11, 5.5, arch_items, font_size=18, color=WHITE)

    # ── SLIDE 5: SCAMWATCH ──────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.8, "SCAMWatch — AI Scam Detection", font_size=36, color=RED, bold=True)

    scam_items = [
        "7 scam types: Digital Arrest, Fake KYC, Fake Investment, Fake Job, Fake Lottery, Impersonation, Romance",
        "5-node LangGraph pipeline: Classify → LLM Analysis → Parse → Score → Store",
        "Risk scoring engine with keyword, urgency, and authority detection",
        "Citizen Alert with emergency contacts (1930, cybercrime.gov.in)",
        "File Complaint button → cybercrime.gov.in (NCRB portal)",
        "Number Block button → sancharsaathi.gov.in (TRAI/Chakshu)",
        "Multi-language: English, Hindi, Tamil, Bengali, Telugu",
        "Accuracy: 85% | False positive rate: 0%",
    ]
    add_bullet_slide(slide, 0.8, 1.5, 11, 5.5, scam_items, font_size=16, color=WHITE)

    # ── SLIDE 6: CURRENCYGUARD ──────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.8, "CURRENCYGuard — Currency Authentication", font_size=36, color=YELLOW, bold=True)

    curr_items = [
        "7 OpenCV security feature checks",
        "Play money / obvious fake detection",
        "Supported: ₹50, ₹100, ₹200, ₹500, ₹2000",
        "Checks: Aspect ratio, Color distribution, Security thread, Serial number, Watermark, Print sharpness",
        "Verdicts: GENUINE, SUSPECT, COUNTERFEIT, INCONCLUSIVE",
        "PDF authenticity reports with feature breakdown",
        "LLM synthesis for expert reasoning",
    ]
    add_bullet_slide(slide, 0.8, 1.5, 11, 5.5, curr_items, font_size=18, color=WHITE)

    # ── SLIDE 7: FRAUDGRAPH ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.8, "FRAUDGraph — Fraud Network Intelligence", font_size=36, color=PURPLE, bold=True)

    fraud_items = [
        "Entity extraction: Phone, Account, Device, Victim, Location",
        "LLM-powered entity extraction from victim statements",
        "Graph analysis: Connected components + betweenness centrality",
        "Interactive pyvis network visualization (dark theme)",
        "Court-admissible evidence kit (ZIP: PDF + Graph HTML + CSV + JSON + Manifest)",
        "Legal disclaimers: IT Act 2000, DPDP Act 2023",
        "Neo4j graph persistence with in-memory fallback",
    ]
    add_bullet_slide(slide, 0.8, 1.5, 11, 5.5, fraud_items, font_size=18, color=WHITE)

    # ── SLIDE 8: GEOINTEL + DASHBOARD ───────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.8, "GeoIntel + Intelligence Dashboard", font_size=36, color=GREEN, bold=True)

    geo_items = [
        "15 NCRB/RBI/MHA reported fraud hotspots across India",
        "Pydeck interactive heatmap with dark theme",
        "Live Threat Monitor with timeline chart",
        "Cross-module correlations (phone/account matching)",
        "Real-time activity feed across all modules",
        "Module health checks and status monitoring",
    ]
    add_bullet_slide(slide, 0.8, 1.5, 11, 5.5, geo_items, font_size=18, color=WHITE)

    # ── SLIDE 9: CITIZEN FRAUD SHIELD ───────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.8, "Citizen Fraud Shield", font_size=36, color=ACCENT_BLUE, bold=True)

    shield_items = [
        "WhatsApp-style chat interface for scam detection",
        "Real-time analysis via SCAMWatch backend",
        "Sample scam messages to try (Digital Arrest, Fake KYC, Fake Investment, Lottery)",
        "Emergency contacts: 1930, cybercrime.gov.in, 181, 112, 14448",
        "File Complaint → NCRB portal",
        "Number Block → Sanchar Saathi / TRAI",
        "Multi-language support for Indian citizens",
    ]
    add_bullet_slide(slide, 0.8, 1.5, 11, 5.5, shield_items, font_size=18, color=WHITE)

    # ── SLIDE 10: INNOVATION HIGHLIGHTS ─────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.8, "Innovation Highlights", font_size=36, color=WHITE, bold=True)

    innov_items = [
        "Cross-module intelligence sharing (unique moat)",
        "Phone in SCAMWatch + FRAUDGraph = correlation alert",
        "Hybrid rule+LLM approach (fast AND intelligent)",
        "Zero-cost pre-classifier before LLM invocation",
        "Play money detection (cost-efficient early exit)",
        "Geospatial hotspot mapping with 15 real locations",
        "Court-admissible evidence kit packaging",
        "85% accuracy with 0 false positives",
    ]
    add_bullet_slide(slide, 0.8, 1.5, 11, 5.5, innov_items, font_size=18, color=WHITE)

    # ── SLIDE 11: IMPACT & SCALABILITY ──────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.8, "Impact & Scalability", font_size=36, color=WHITE, bold=True)

    impact_items = [
        "Citizens: Real-time scam detection before money transfer",
        "Bank tellers: Instant currency authentication",
        "Law enforcement: Fraud ring mapping + evidence packages",
        "Police: Geographic patrol prioritisation",
        "Docker Compose ready for deployment",
        "Cloud-ready: Neo4j Aura, Groq API, ChromaDB",
        "Modular architecture — new modules without touching existing code",
    ]
    add_bullet_slide(slide, 0.8, 1.5, 11, 5.5, impact_items, font_size=18, color=WHITE)

    # ── SLIDE 12: DEMO & THANK YOU ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 1, 1.0, 11, 1, "Demo & Next Steps", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    demo_items = [
        "Live demo: SCAMWatch → CURRENCYGuard → FRAUDGraph → GeoIntel",
        "Future: WhatsApp Business API, IVR integration, 12 languages",
        "Future: Real-time alerting, mobile app deployment",
        "",
        "Thank you!",
        "GitHub: github.com/Sushit-prog/SENTINEL",
    ]
    add_bullet_slide(slide, 0.8, 2.5, 11, 4, demo_items, font_size=18, color=WHITE)

    # Save
    output_path = "SENTINEL_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
